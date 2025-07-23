#!/usr/bin/env python
"""Robust extractor + captioner with MIME logging and rich progress."""
from logconf import logging
from pathlib import Path
from itertools import chain
from rich.progress import Progress
import mimetypes, uuid, json, traceback, torch

from unstructured.partition.auto import partition
from unstructured.documents.elements import Element
from transformers import BlipProcessor, BlipForConditionalGeneration
from pptx import Presentation
from PIL import Image

BASE      = Path(__file__).resolve().parent.parent
RAW       = BASE / "raw"
RAW_IMG   = BASE / "raw_imgs"
CLEAN     = BASE / "clean"
RAW_IMG.mkdir(exist_ok=True)
CLEAN.mkdir(exist_ok=True)

# 🔍 Captioning model (BLIP base)
proc  = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
model = BlipForConditionalGeneration.from_pretrained(
    "Salesforce/blip-image-captioning-base", torch_dtype=torch.float16
).to("mps")

def caption(img_path: Path) -> str:
    try:
        inputs = proc(images=Image.open(img_path), return_tensors="pt").to("mps")
        ids = model.generate(**inputs, max_new_tokens=25)
        return proc.decode(ids[0], skip_special_tokens=True)
    except Exception:
        logging.exception(f"⚠️ Caption failed for {img_path}")
        return "image"

def textify(el: Element) -> str:
    return getattr(el, "to_markdown", lambda: el.text)()

docs = chain(
    RAW.rglob("*.pptx"),
    RAW.rglob("*.docx"),
    RAW.rglob("*.pdf"),
    RAW.rglob("*.xlsx"),
    RAW.rglob("*.vsdx"),
)

with Progress() as bar:
    for fp in bar.track(list(docs), description="📁 Extracting and captioning"):
        if fp.name.startswith("~$"):
            logging.warning(f"⏭️ Skipped temporary/system file: {fp.name}")
            continue

        try:
            mimetype, _ = mimetypes.guess_type(fp)
            logging.info(f"📄 Processing {fp.name} ({mimetype})")

            doc_id = uuid.uuid4().hex
            try:
                els = partition(str(fp))
                md = "\n".join(textify(e) for e in els)
            except Exception:
                logging.error(f"❌ Partition failed for {fp.name}\n{traceback.format_exc()}")
                continue

            # 🎞️ Extract slide images + caption for PPTX
            if fp.suffix.lower() == ".pptx":
                try:
                    pres = Presentation(fp)
                    for idx, slide in enumerate(pres.slides):
                        for shp in slide.shapes:
                            if shp.shape_type == 13:  # Picture
                                img = RAW_IMG / f"{doc_id}_{idx}.png"
                                with open(img, "wb") as f:
                                    f.write(shp.image.blob)
                                md += f"\n\n![{caption(img)}]({img})"
                except Exception:
                    logging.warning(f"⚠️ Failed to extract images from {fp.name}")

            # 🧾 Save extracted content
            (CLEAN / f"{doc_id}.json").write_text(json.dumps({
                "id": doc_id, "title": fp.stem, "body": md, "source": str(fp)
            }))
            logging.info(f"📝 Extracted {fp.name}")

        except Exception:
            logging.error(f"❌ Skipped {fp.name}\n{traceback.format_exc()}")